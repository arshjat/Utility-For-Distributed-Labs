
# IMPORT PACKAGES
from pptx import Presentation
import os

def extract_text():
    """
    This function extracts data from ppt slides that are uploaded by the staff.
    Input :        Power Point Presentation / Slides
    Output:        Extracted Text
    Requirements:  None
    
    """


    # PATH
    cur_directory = os.path.dirname(os.path.abspath( __file__ ))
    files = [x for x in os.listdir(cur_directory) if x.endswith(".pptx")]
    
    # EXTRACT
    text_runs = []
    for each in files:
        prs = Presentation(each)
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_runs.append(run.text)
    return text_runs